import os
import io
import re
import random
import traceback
import logging
from contextlib import asynccontextmanager

import google.generativeai as genai
import httpx
import nltk
import numpy as np
import PyPDF2
import uvicorn
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from html.parser import HTMLParser
from pptx import Presentation
from pydantic import BaseModel
from rank_bm25 import BM25Okapi
from sumy.nlp.tokenizers import Tokenizer
from sumy.parsers.plaintext import PlaintextParser
from sumy.summarizers.text_rank import TextRankSummarizer
from supabase import Client, create_client
from typing import List, Optional

_NLTK_PACKAGES = ("punkt", "punkt_tab", "stopwords")


def _gemini_models():
    return (
        os.getenv("GEMINI_MODEL", "gemini-2.5-flash"),
        "gemini-2.5-flash",
        "gemini-2.0-flash-lite",
        "gemini-2.0-flash",
    )


def _nltk_search_paths():
    paths = []
    for candidate in (
        os.environ.get("NLTK_DATA"),
        "/usr/local/share/nltk_data",
        "/tmp/nltk_data",
    ):
        if candidate and candidate not in paths:
            paths.append(candidate)
    return paths


def _ensure_nltk_data():
    for download_dir in _nltk_search_paths():
        os.makedirs(download_dir, exist_ok=True)
        if download_dir not in nltk.data.path:
            nltk.data.path.insert(0, download_dir)

    for package in _NLTK_PACKAGES:
        try:
            if "punkt" in package:
                nltk.data.find(f"tokenizers/{package}")
            else:
                nltk.data.find(f"corpora/{package}")
        except LookupError:
            download_dir = _nltk_search_paths()[-1]
            logging.info("Downloading NLTK package %s to %s", package, download_dir)
            nltk.download(package, download_dir=download_dir, quiet=True)


logging.basicConfig(level=logging.INFO)
load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")


@asynccontextmanager
async def lifespan(_: FastAPI):
    _ensure_nltk_data()
    yield


app = FastAPI(title="Adaptivise AI Engine", lifespan=lifespan)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

SUPABASE_URL = os.getenv("SUPABASE_URL", "https://nnqafyfydbpspywuxhtk.supabase.co")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY") or os.getenv("SUPABASE_KEY", "")
supabase: Optional[Client] = None
if SUPABASE_KEY:
    supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

_STOPWORDS = {
    "about", "after", "also", "been", "being", "between", "could", "each",
    "from", "have", "into", "more", "other", "such", "than", "that", "their",
    "them", "then", "there", "these", "they", "this", "those", "through",
    "under", "very", "were", "what", "when", "where", "which", "while",
    "with", "would", "your", "the", "and", "a", "to", "of", "in", "i", "is",
    "it", "on", "you", "for", "but", "are", "be", "at", "or", "as", "was",
    "so", "if", "out", "not", "by", "an"
}

def generate_distractors(answer, all_sentences, count=3):
    pool = [
        s for s in all_sentences
        if s.strip().lower() != answer.strip().lower() and len(s) > 10
    ]
    if len(pool) <= count:
        return pool
    return random.sample(pool, count)

def truncate_text(text: str, limit: int = 65) -> str:
    """Aggressively truncates text for mobile UI buttons to prevent overflow."""
    # Remove any weird line breaks that might have been extracted
    cleaned = text.replace('\n', ' ').strip()
    if len(cleaned) > limit:
        # Cut it off cleanly without splitting a word in half
        return cleaned[:limit].rsplit(' ', 1)[0] + "..."
    return cleaned

def _fallback_quiz():
    """Provides a safe fallback question if the text is unreadable or too short."""
    return [{
        "question": "Review the uploaded material to understand the core concepts.",
        "answer": "Understood the material.",
        "options": [
            "Understood the material.",
            "I need to read it again.",
            "The material was too complex.",
            "None of the above."
        ]
    }]

def generate_semantic_quiz(text: str) -> list:
    try:
        cleaned = " ".join(text.split())
        if len(cleaned) < 15:
            return _fallback_quiz()

        sentences = [
            s.strip() for s in re.split(r"(?<=[.!?]) +", cleaned)
            if len(s.strip()) > 12
        ]

        if len(sentences) == 0:
            sentences = [cleaned[i:i+120] for i in range(0, len(cleaned), 120) if len(cleaned[i:i+120]) > 12]

        # DYNAMIC QUESTION COUNT: based on document length (max 10 questions)
        max_questions = max(3, min(10, len(sentences) // 3))

        if len(sentences) == 1:
            sentence = truncate_text(sentences[0], 110)
            return [{
                "question": "Which statement best matches the material?",
                "answer": sentence,
                "options": [
                    sentence,
                    "This topic is unrelated to the uploaded material.",
                    "The material contradicts this statement.",
                    "None of the above."
                ],
            }]

        tokenized = [s.lower().split() for s in sentences if s.strip()]
        if not tokenized:
            return _fallback_quiz()

        bm25 = BM25Okapi(tokenized)
        scores = list(bm25.get_scores(cleaned.lower().split()))
        ranked = sorted(range(len(sentences)), key=lambda i: scores[i], reverse=True)
        selected = ranked[: min(max_questions, len(ranked))]

        quiz = []
        for idx in selected:
            sentence = sentences[idx]
            keywords = [
                w for w in re.findall(r"\b[A-Za-z]{4,}\b", sentence)
                if w.lower() not in _STOPWORDS
            ]

            if keywords:
                key_term = max(keywords, key=len)
                question = f"According to the material, which statement is correct about '{key_term}'?"
            else:
                question = "According to the material, which statement is correct?"

            raw_answer = sentence
            distractor_pool = [
                sentences[i] for i in ranked
                if i != idx and sentences[i].strip().lower() != raw_answer.strip().lower()
            ]
            
            base_distractors = generate_distractors(raw_answer, distractor_pool, count=3)
            distractors = []
            for d in base_distractors:
                # Add a prefix to invalidate the statement in the context of the question
                distractors.append(f"Incorrectly states: {d[:80].lower()}...")
            
            if len(distractors) < 3:
                distractors.extend([
                    "This statement is unrelated to the uploaded material.",
                    "The material contradicts this statement.",
                    "None of the above apply.",
                ])
                
            # TRUNCATE to fix Kinesthetic UI overflow
            answer = truncate_text(raw_answer, 110)
            distractors = [truncate_text(d, 110) for d in distractors]
                
            options = list(dict.fromkeys([answer] + distractors))[:4]
            random.shuffle(options)

            quiz.append({
                "question": question,
                "answer": answer,
                "options": options,
            })

        return quiz if quiz else _fallback_quiz()

    except Exception as e:
        print(f"======== CRITICAL ERROR IN SEMANTIC QUIZ: {e} ========")
        traceback.print_exc()
        return _fallback_quiz()

def extract_text_from_pptx(content: bytes) -> str:
    text_runs = []
    prs = Presentation(io.BytesIO(content))
    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text.strip())
    return "\n".join(text_runs)

def extract_text_from_file(content: bytes, filename: str) -> str:
    text = ""
    file_extension = filename.split('.')[-1].lower()

    if file_extension == 'pdf':
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    elif file_extension == 'docx':
        doc = Document(io.BytesIO(content))
        text = "\n".join([para.text for para in doc.paragraphs])
    elif file_extension == 'pptx':
        text = extract_text_from_pptx(content)
    elif file_extension == 'txt':
        text = content.decode("utf-8", errors="replace")

    return text

class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._parts = []

    def handle_data(self, data):
        text = data.strip()
        if text:
            self._parts.append(text)

    def get_text(self):
        return " ".join(self._parts)

def extract_text_from_html(content: bytes) -> str:
    html = content.decode("utf-8", errors="replace")
    parser = _HTMLTextExtractor()
    parser.feed(html)
    return parser.get_text()

def extract_text_from_url(url: str) -> str:
    with httpx.Client(timeout=20.0, follow_redirects=True) as client:
        response = client.get(url)
        response.raise_for_status()
        content_type = response.headers.get("content-type", "").lower()
        if "html" in content_type or url.lower().endswith((".html", ".htm")):
            return extract_text_from_html(response.content)
        return response.text

def simple_summarize(text_input: str):
    sentences = [s.strip() for s in re.split(r'(?<=[.!?]) +', text_input) if s.strip()]
    
    # DYNAMIC LENGTH: Extracts up to 15 sentences depending on document size (fixes short audio)
    num_sentences = max(4, min(15, len(sentences) // 3))
    if len(sentences) <= num_sentences:
        return " ".join(sentences)
            
    words = [w.lower() for w in re.findall(r'\w+', text_input.lower()) if len(w) > 3]
    freq = {}
    for w in words:
        freq[w] = freq.get(w, 0) + 1
        
    scores = []
    for s in sentences:
        s_words = [w for w in re.findall(r'\w+', s.lower()) if len(w) > 3]
        if len(s_words) > 0:
            # FIX: Normalize by sentence length to prevent bias towards long sentences
            score = sum(freq.get(w, 0) for w in s_words) / len(s_words)
        else:
            score = 0
        scores.append(score)
        
    indices = np.argsort(scores)[-num_sentences:]
    indices.sort()
    return " ".join([sentences[i] for i in indices])

def _format_summary_prompt(raw_summary: str) -> str:
    return f"""You are an expert educational editor designing content for a mobile app.
Reformat this raw summary for excellent mobile readability using clean Markdown.

STRICT RULES:
1. Start with a single # Heading for the main title.
2. Convert the text into short, highly scannable bullet points.
3. Add a blank line between every bullet point.
4. **Bold** 1 or 2 critical keywords per bullet point.
5. Do not add information that is not in the raw summary.
6. Output valid Markdown only.

Raw summary:
{raw_summary}"""


def _local_format_summary(raw_summary: str) -> str:
    sentences = [s.strip() for s in re.split(r"(?<=[.!?]) +", raw_summary) if s.strip()]
    if not sentences:
        return raw_summary

    title = sentences[0]
    if len(title) > 80:
        title = title[:77] + "..."

    lines = [f"# {title}", ""]
    for sentence in sentences:
        lines.extend([f"- {sentence}", ""])
    return "\n".join(lines).strip()


def format_summary_with_gemini(raw_summary: str) -> str:
    cleaned = " ".join(raw_summary.split()).strip()
    if not cleaned:
        return cleaned
    if not GEMINI_API_KEY:
        return _local_format_summary(cleaned)

    genai.configure(api_key=GEMINI_API_KEY)
    prompt = _format_summary_prompt(cleaned)

    for model_name in dict.fromkeys(_gemini_models()):
        try:
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(prompt)
            formatted = (response.text or "").strip()
            if formatted:
                return formatted
        except Exception as exc:
            logging.warning("Gemini model %s failed: %s", model_name, exc)

    return _local_format_summary(cleaned)

class NotesRequest(BaseModel):
    text: str

class VarkScores(BaseModel):
    scores: List[int]

@app.get("/health")
def health_check():
    nltk_ok = True
    try:
        _ensure_nltk_data()
        nltk.data.find("tokenizers/punkt_tab")
    except Exception:
        nltk_ok = False
    return {"status": "ok", "nltk_ready": nltk_ok}

@app.post("/ai/process-note")
async def process_note(
    file: UploadFile = File(...),
    user_id: str = Form(...),
    storage_path: str = Form(...),
    folder_id: str = Form(None),
):
    logging.info("Processing started for file: %s", file.filename)
    try:
        content = await file.read()
        raw_text = extract_text_from_file(content, file.filename)

        if not raw_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from file.")

        cleaned_text = " ".join(raw_text.split())
        summary_text = simple_summarize(cleaned_text)
        quiz_content = generate_semantic_quiz(cleaned_text)

        response = {
            "status": "success",
            "user_id": user_id,
            "storage_path": storage_path,
            "file_name": file.filename,
            "raw_text": cleaned_text,
            "summary": summary_text,
            "quiz_content": quiz_content,
        }
        if folder_id:
            response["folder_id"] = folder_id

        return response
    except HTTPException:
        raise
    except Exception as e:
        logging.error("CRASHED during processing: %s", traceback.format_exc())
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/ai/process-url")
async def process_url(
    url: str = Form(...),
    user_id: str = Form(...),
    storage_path: str = Form(...),
    folder_id: str = Form(None),
):
    logging.info("Processing URL: %s", url)
    try:
        raw_text = extract_text_from_url(url.strip())
        if not raw_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from URL.")

        cleaned_text = " ".join(raw_text.split())
        summary_text = simple_summarize(cleaned_text)
        quiz_content = generate_semantic_quiz(cleaned_text)

        response = {
            "status": "success",
            "user_id": user_id,
            "storage_path": storage_path,
            "file_name": url.split("/")[-1] or "web-page.html",
            "raw_text": cleaned_text,
            "summary": summary_text,
            "quiz_content": quiz_content,
        }
        if folder_id:
            response["folder_id"] = folder_id
        return response
    except HTTPException:
        raise
    except Exception as e:
        logging.error("URL processing failed: %s", traceback.format_exc())
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/ai/format-summary")
async def format_summary_endpoint(request: NotesRequest):
    try:
        return {"formatted_summary": format_summary_with_gemini(request.text)}
    except Exception as e:
        logging.error("Summary formatting failed: %s", traceback.format_exc())
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/ai/summarize")
async def summarize_notes(request: NotesRequest):
    try:
        _ensure_nltk_data()
        cleaned_text = " ".join(request.text.split())
        parser = PlaintextParser.from_string(cleaned_text, Tokenizer("english"))
        sentences = parser.document.sentences
        
        sentence_list = [str(s) for s in sentences]
        tokenized_corpus = [s.split(" ") for s in sentence_list]
        bm25 = BM25Okapi(tokenized_corpus)
        
        bm25_raw_scores = bm25.get_scores(cleaned_text.split(" "))
        max_bm25 = max(bm25_raw_scores) if max(bm25_raw_scores) > 0 else 1
        norm_bm25 = [score / max_bm25 for score in bm25_raw_scores]
        
        summarizer = TextRankSummarizer()
        tr_scores_dict = summarizer.rate_sentences(parser.document)
        max_tr = max(tr_scores_dict.values()) if tr_scores_dict.values() else 1

        hybrid_scored_list = []
        for i, sentence in enumerate(sentences):
            tr_score = tr_scores_dict.get(sentence, 0) / max_tr
            bm25_score = norm_bm25[i]
            final_score = (tr_score * 0.5) + (bm25_score * 0.5)
            hybrid_scored_list.append((sentence, final_score))
        
        hybrid_scored_list.sort(key=lambda x: x[1], reverse=True)
        
        # DYNAMIC LENGTH for longer audio notes
        summary_length = max(4, min(15, int(len(sentences) * 0.35)))
        top_sentences = hybrid_scored_list[:summary_length]
        top_sentences.sort(key=lambda x: sentences.index(x[0]))
        summary_text = " ".join([str(s[0]).strip() for s in top_sentences])

        return {
            "summary": summary_text,
            "bm25_relevance_peak": max(bm25_raw_scores) if len(bm25_raw_scores) > 0 else 0,
            "keywords": [] 
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/ai/generate-quiz")
async def generate_quiz_endpoint(request: NotesRequest):
    try:
        cleaned_text = " ".join(request.text.split())
        return {"quiz_content": generate_semantic_quiz(cleaned_text)}
    except Exception as e:
        print("======== FASTAPI ROUTE ERROR ========")
        traceback.print_exc()
        return {"quiz_content": _fallback_quiz()}

@app.post("/ai/classify-vark")
async def classify_vark(request: VarkScores):
    scores = request.scores
    
    # Run the actual Hybrid Weighted K-Means algorithm
    detected_styles = hybrid_weighted_kmeans(scores)
    
    primary_style = " & ".join(detected_styles) if len(detected_styles) > 1 else detected_styles[0]
    clean_primary_style = primary_style.encode('ascii', 'ignore').decode('ascii')
    
    return {
        "learning_style": detected_styles,
        "formatted_style": clean_primary_style
    }

def hybrid_weighted_kmeans(new_user_scores, k=4, max_iters=100):
    """
    A true Hybrid Weighted K-Means clustering algorithm.
    It groups the new user alongside a historical dataset using weighted Euclidean distances.
    """
    # 1. Simulated Historical Dataset of previous student scores 
    # (This allows the K-Means algorithm to actually form clusters)
    historical_data = [
        [15, 2, 1, 0], [14, 3, 2, 1], [16, 0, 1, 2], # Mostly Visual
        [1, 16, 2, 0], [2, 14, 1, 3], [0, 15, 3, 1], # Mostly Auditory
        [2, 1, 15, 2], [0, 2, 16, 1], [3, 0, 14, 2], # Mostly Read/Write
        [1, 2, 0, 16], [2, 1, 3, 14], [0, 3, 1, 15], # Mostly Kinesthetic
    ]
    
    # Append the new user to the dataset
    dataset = np.array(historical_data + [new_user_scores], dtype=float)
    
    # 2. Apply Weights (The "Weighted" part of HWKM)
    # Applying slight biases to Auditory and Kinesthetic as per your original logic
    weights = np.array([1.0, 1.1, 1.0, 1.1])
    weighted_dataset = dataset * weights
    
    # 3. Smart Initialization (The "Hybrid" part)
    # Instead of choosing random starting points (which causes instability), 
    # we use deterministic extreme points to seed the initial centroids.
    centroids = np.array([
        [16.0, 0.0, 0.0, 0.0] * weights,
        [0.0, 16.0, 0.0, 0.0] * weights,
        [0.0, 0.0, 16.0, 0.0] * weights,
        [0.0, 0.0, 0.0, 16.0] * weights
    ])
    
    labels = np.zeros(len(weighted_dataset))
    
    # 4. The Iterative K-Means Loop
    for _ in range(max_iters):
        # Assign each student to the nearest centroid using Euclidean distance
        for i, point in enumerate(weighted_dataset):
            distances = [np.sqrt(np.sum((point - c)**2)) for c in centroids]
            labels[i] = np.argmin(distances)
            
        # Recalculate centroids based on the mean of the assigned points
        new_centroids = np.copy(centroids)
        for cluster_idx in range(k):
            cluster_points = weighted_dataset[labels == cluster_idx]
            if len(cluster_points) > 0:
                new_centroids[cluster_idx] = np.mean(cluster_points, axis=0)
                
        # Check for convergence (if centroids stop moving, the algorithm is done)
        if np.allclose(centroids, new_centroids):
            break
        centroids = new_centroids

    # 5. Extract the assigned cluster for our NEW USER (who is the very last item in the array)
    new_user_cluster_idx = int(labels[-1])
    assigned_centroid = centroids[new_user_cluster_idx] / weights # Remove weights to read it normally
    
    style_names = ['Visual', 'Auditory', 'Read/Write', 'Kinesthetic']
    
    # Determine what this cluster represents based on its highest average score
    dominant_feature_idx = np.argmax(assigned_centroid)
    detected_style = [style_names[dominant_feature_idx]]
    
    return detected_style